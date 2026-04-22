"""
Advanced Document Parser
Xử lý PDF/DOCX với khả năng trích xuất:
- Text (OCR cho scanned PDF)
- Images (embedded images) với position-aware OCR
- Tables (structured data)
- Metadata (author, title, etc.)
"""
import logging
import io
import base64
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import List, Dict, Any, Optional
from dataclasses import dataclass

logger = logging.getLogger(__name__)


@dataclass
class ParsedContent:
    """Kết quả parse document"""
    text: str
    images: List[Dict[str, Any]]  # [{"page": 1, "data": base64, "caption": "..."}]
    tables: List[Dict[str, Any]]  # [{"page": 1, "data": [[...]], "caption": "..."}]
    metadata: Dict[str, Any]
    total_pages: int


class AdvancedDocumentParser:
    """Parser nâng cao cho PDF và DOCX"""

    MAX_EMBEDDED_IMAGE_OCR = 10
    MIN_EMBEDDED_IMAGE_SIZE_BYTES = 10_000
    
    def __init__(self):
        self._pdf_parser = None
        self._docx_parser = None

    def _should_ocr_embedded_image(self, image_bytes: bytes, ocr_count: int) -> bool:
        """Only OCR meaningful embedded images and stop after a safe cap."""
        return (
            len(image_bytes) >= self.MIN_EMBEDDED_IMAGE_SIZE_BYTES
            and ocr_count < self.MAX_EMBEDDED_IMAGE_OCR
        )
    
    async def parse_pdf(self, file_path: Path) -> ParsedContent:
        """
        Parse PDF với khả năng:
        - OCR cho scanned PDF
        - Trích xuất hình ảnh
        - Trích xuất bảng biểu
        - Metadata
        """
        try:
            import fitz  # PyMuPDF
            import pdfplumber
            from PIL import Image
            
            logger.info(f"Parsing PDF: {file_path}")
            
            # Open with PyMuPDF for images and metadata
            pdf_doc = fitz.open(str(file_path))
            total_pages = len(pdf_doc)
            
            # Extract metadata
            metadata = {
                "title": pdf_doc.metadata.get("title", ""),
                "author": pdf_doc.metadata.get("author", ""),
                "subject": pdf_doc.metadata.get("subject", ""),
                "creator": pdf_doc.metadata.get("creator", ""),
                "producer": pdf_doc.metadata.get("producer", ""),
                "creation_date": pdf_doc.metadata.get("creationDate", ""),
                "total_pages": total_pages
            }
            
            text_parts = []
            images = []
            tables = []
            ocr_count = 0
            
            # Open with pdfplumber for text and tables
            with pdfplumber.open(str(file_path)) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    logger.info(f"Processing page {page_num}/{total_pages}")
                    
                    # Extract text
                    page_text = page.extract_text()
                    
                    # If no text, try OCR
                    if not page_text or len(page_text.strip()) < 20:
                        logger.info(f"Page {page_num} has little text, trying OCR...")
                        page_text = await self._ocr_page(pdf_doc[page_num - 1])
                    
                    if page_text:
                        text_parts.append(f"# Page {page_num}\n\n{page_text}\n")
                    
                    # Extract tables
                    page_tables = page.extract_tables()
                    if page_tables:
                        for table_idx, table in enumerate(page_tables):
                            if table and len(table) > 0:
                                tables.append({
                                    "page": page_num,
                                    "table_index": table_idx,
                                    "data": table,
                                    "caption": f"Table {table_idx + 1} on page {page_num}"
                                })
                                
                                # Add table to text
                                table_text = self._format_table_as_text(table)
                                text_parts.append(f"\n[Table {table_idx + 1}]\n{table_text}\n")
            
            # Extract images with PyMuPDF
            # OCR only meaningful embedded images and cap OCR count per file.
            ocr_count = 0
            
            for page_num in range(total_pages):
                page = pdf_doc[page_num]
                image_list = page.get_images()
                
                for img_idx, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        base_image = pdf_doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        
                        # OCR only meaningful images (large enough + under cap)
                        image_text = ""
                        if self._should_ocr_embedded_image(image_bytes, ocr_count):
                            logger.info(f"OCR image {img_idx+1} on page {page_num+1} ({len(image_bytes)//1024}KB)")
                            image_text = await self._ocr_image(image_bytes)
                            ocr_count += 1
                        
                        images.append({
                            "page": page_num + 1,
                            "image_index": img_idx,
                            "format": image_ext,
                            "data": "",  # Don't store base64 to save memory
                            "caption": f"Image {img_idx + 1} on page {page_num + 1}",
                            "extracted_text": image_text
                        })
                        
                        # Add image text to main text
                        if image_text:
                            text_parts.append(f"\n[Image {img_idx + 1} - Page {page_num + 1}]\n{image_text}\n")
                        
                    except Exception as e:
                        logger.warning(f"Error extracting image {img_idx} from page {page_num}: {e}")
            
            if ocr_count > 0:
                logger.info(f"OCR'd {ocr_count} images (skipped {len(images) - ocr_count} small/excess)")
            
            pdf_doc.close()
            
            full_text = "\n".join(text_parts)
            
            if not full_text or len(full_text.strip()) < 50:
                raise ValueError("Không thể trích xuất nội dung từ PDF")
            
            logger.info(f"✅ Parsed PDF: {len(full_text)} chars, {len(images)} images, {len(tables)} tables")
            
            return ParsedContent(
                text=full_text,
                images=images,
                tables=tables,
                metadata=metadata,
                total_pages=total_pages
            )
            
        except ImportError as e:
            logger.error(f"Missing dependency: {e}")
            raise ValueError(f"Thiếu thư viện: {str(e)}. Cần rebuild Docker image.")
        except Exception as e:
            logger.error(f"PDF parsing error: {e}")
            raise ValueError(f"Lỗi parse PDF: {str(e)}")
    
    async def parse_docx(self, file_path: Path) -> ParsedContent:
        """
        Parse DOCX với khả năng:
        - Trích xuất text với formatting
        - Trích xuất hình ảnh
        - Trích xuất bảng biểu
        - Metadata
        """
        try:
            import docx
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
            from docx.table import _Cell, Table
            from docx.text.paragraph import Paragraph
            from PIL import Image
            
            logger.info(f"Parsing DOCX: {file_path}")
            
            doc = docx.Document(str(file_path))
            
            # Extract metadata
            core_props = doc.core_properties
            metadata = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "keywords": core_props.keywords or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "total_pages": len(doc.sections)
            }
            
            text_parts = []
            images = []
            tables = []
            
            # Process document elements in order
            for element in doc.element.body:
                # Paragraph
                if isinstance(element, CT_P):
                    para = Paragraph(element, doc)
                    if para.text.strip():
                        # Check if heading
                        if para.style.name.startswith('Heading'):
                            level = para.style.name.replace('Heading ', '')
                            if level.isdigit():
                                text_parts.append(f"{'#' * int(level)} {para.text}\n")
                            else:
                                text_parts.append(f"# {para.text}\n")
                        else:
                            text_parts.append(para.text)
                
                # Table
                elif isinstance(element, CT_Tbl):
                    table = Table(element, doc)
                    table_data = []
                    
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    
                    if table_data:
                        tables.append({
                            "page": 0,  # DOCX doesn't have page numbers easily
                            "table_index": len(tables),
                            "data": table_data,
                            "caption": f"Table {len(tables) + 1}"
                        })
                        
                        # Add table to text
                        table_text = self._format_table_as_text(table_data)
                        text_parts.append(f"\n[Table {len(tables)}]\n{table_text}\n")
            
            # ── Position-aware DOCX image OCR ──
            # Parse the DOCX XML to find images in their correct paragraph positions
            # and insert OCR text inline (inspired by RAG-Anything implementation)
            logger.info("[AdvancedParser] DOCX inline image OCR: parsing XML structure...")
            
            try:
                with zipfile.ZipFile(str(file_path), "r") as docx_zip:
                    # 1. Build rId → media filename map from relationships
                    rid_to_media = {}
                    try:
                        rels_xml = docx_zip.read("word/_rels/document.xml.rels")
                        rels_root = ET.fromstring(rels_xml)
                        for rel in rels_root:
                            target = rel.get("Target", "")
                            if target.startswith("media/"):
                                rid_to_media[rel.get("Id")] = f"word/{target}"
                    except Exception:
                        pass
                    
                    if not rid_to_media:
                        logger.info("[AdvancedParser] No image relationships found in DOCX")
                    else:
                        logger.info(f"[AdvancedParser] Found {len(rid_to_media)} image refs in DOCX")
                        
                        # 2. Pre-OCR all unique media files
                        media_ocr_cache = {}  # media_path → ocr text
                        ocr_count = 0
                        for rid, media_path in rid_to_media.items():
                            if media_path in media_ocr_cache:
                                continue
                            try:
                                img_data = docx_zip.read(media_path)
                                img_b64 = base64.b64encode(img_data).decode('utf-8')
                                ocr_text = ""
                                if self._should_ocr_embedded_image(img_data, ocr_count):
                                    logger.info(
                                        f"[AdvancedParser] DOCX OCR {media_path} ({len(img_data)//1024}KB)"
                                    )
                                    ocr_text = await self._ocr_image(img_data)
                                    ocr_count += 1
                                media_ocr_cache[media_path] = ocr_text
                                
                                # Store image metadata
                                images.append({
                                    "page": 0,
                                    "image_index": len(images),
                                    "format": media_path.split('.')[-1],
                                    "data": img_b64,
                                    "caption": f"Image {len(images) + 1}",
                                    "extracted_text": ocr_text,
                                    "rid": rid
                                })
                                
                            except Exception as e:
                                logger.warning(f"Failed to OCR {media_path}: {e}")
                                media_ocr_cache[media_path] = ""
                        if ocr_count > 0:
                            logger.info(
                                f"[AdvancedParser] DOCX OCR'd {ocr_count} images "
                                f"(skipped {max(0, len(media_ocr_cache) - ocr_count)} small/excess)"
                            )
                        
                        # 3. Walk document.xml paragraphs in order,
                        #    inserting image OCR text at the correct position
                        doc_xml = docx_zip.read("word/document.xml")
                        doc_root = ET.fromstring(doc_xml)
                        
                        # XML namespaces used in OOXML
                        ns = {
                            "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
                            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
                            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                            "wp": "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
                            "pic": "http://schemas.openxmlformats.org/drawingml/2006/picture",
                        }
                        
                        # Find all paragraphs inside <w:body>
                        body = doc_root.find(".//w:body", ns)
                        if body is None:
                            body = doc_root
                        
                        inline_parts = []  # list of text chunks in reading order
                        img_count = 0
                        
                        for elem in body:
                            tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
                            
                            if tag == "p":
                                # Extract plain text from this paragraph
                                para_texts = []
                                for t_elem in elem.iter(f"{{{ns['w']}}}t"):
                                    if t_elem.text:
                                        para_texts.append(t_elem.text)
                                para_str = "".join(para_texts).strip()
                                
                                # Check if this paragraph contains any drawings
                                drawings = list(elem.iter(f"{{{ns['w']}}}drawing"))
                                img_texts_in_para = []
                                for drawing in drawings:
                                    # Find <a:blip r:embed="rIdX">
                                    for blip in drawing.iter(f"{{{ns['a']}}}blip"):
                                        embed_rid = blip.get(f"{{{ns['r']}}}embed")
                                        if embed_rid and embed_rid in rid_to_media:
                                            media_path = rid_to_media[embed_rid]
                                            ocr_text = media_ocr_cache.get(media_path, "")
                                            if ocr_text.strip():
                                                img_count += 1
                                                img_texts_in_para.append(ocr_text)
                                
                                # Build the paragraph output
                                if para_str:
                                    inline_parts.append(para_str)
                                # Insert image OCR text RIGHT AFTER the paragraph text
                                for it in img_texts_in_para:
                                    inline_parts.append(f"[Image Text: {it}]")
                            
                            elif tag == "tbl":
                                # Tables: extract text from all cells
                                tbl_texts = []
                                for t_elem in elem.iter(f"{{{ns['w']}}}t"):
                                    if t_elem.text:
                                        tbl_texts.append(t_elem.text)
                                if tbl_texts:
                                    inline_parts.append(" | ".join(tbl_texts))
                        
                        if img_count > 0:
                            # Rebuild full_text and markdown_text with inline image text
                            text_parts = inline_parts
                            logger.info(
                                f"[AdvancedParser] DOCX inline OCR: inserted {img_count} "
                                f"image(s) text at correct positions"
                            )
                        else:
                            logger.info("[AdvancedParser] DOCX images had no extractable text")
            
            except Exception as e:
                logger.warning(f"[AdvancedParser] Failed DOCX inline image OCR: {e}, using fallback")
                # Fallback: Extract images without position awareness
                ocr_count = 0
                for rel in doc.part.rels.values():
                    if "image" in rel.target_ref:
                        try:
                            image_data = rel.target_part.blob
                            image_b64 = base64.b64encode(image_data).decode('utf-8')
                            image_text = ""
                            if self._should_ocr_embedded_image(image_data, ocr_count):
                                logger.info(
                                    f"[AdvancedParser] DOCX fallback OCR {rel.target_ref} ({len(image_data)//1024}KB)"
                                )
                                image_text = await self._ocr_image(image_data)
                                ocr_count += 1
                            
                            images.append({
                                "page": 0,
                                "image_index": len(images),
                                "format": rel.target_ref.split('.')[-1],
                                "data": image_b64,
                                "caption": f"Image {len(images) + 1}",
                                "extracted_text": image_text
                            })
                            
                            if image_text:
                                text_parts.append(f"\n[Image {len(images)}]\n{image_text}\n")
                        except Exception as e2:
                            logger.warning(f"Error extracting image: {e2}")
                if ocr_count > 0:
                    logger.info(
                        f"[AdvancedParser] DOCX fallback OCR'd {ocr_count} images"
                    )
            
            full_text = "\n\n".join(text_parts)
            
            if not full_text or len(full_text.strip()) < 50:
                raise ValueError("Không thể trích xuất nội dung từ DOCX")
            
            logger.info(f"✅ Parsed DOCX: {len(full_text)} chars, {len(images)} images, {len(tables)} tables")
            
            return ParsedContent(
                text=full_text,
                images=images,
                tables=tables,
                metadata=metadata,
                total_pages=len(doc.sections)
            )
            
        except ImportError as e:
            logger.error(f"Missing dependency: {e}")
            raise ValueError(f"Thiếu thư viện: {str(e)}. Cần rebuild Docker image.")
        except Exception as e:
            logger.error(f"DOCX parsing error: {e}")
            raise ValueError(f"Lỗi parse DOCX: {str(e)}")
    
    async def _ocr_page(self, page) -> str:
        """OCR a PDF page using pytesseract"""
        try:
            import pytesseract
            from pdf2image import convert_from_path
            from PIL import Image
            
            # Convert page to image
            pix = page.get_pixmap(dpi=300)
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            
            # OCR
            text = pytesseract.image_to_string(img, lang='vie+eng')
            return text.strip()
            
        except Exception as e:
            logger.warning(f"OCR failed: {e}")
            return ""
    
    async def _ocr_image(self, image_bytes: bytes) -> str:
        """OCR an image using pytesseract"""
        try:
            import pytesseract
            from PIL import Image
            
            img = Image.open(io.BytesIO(image_bytes))
            text = pytesseract.image_to_string(img, lang='vie+eng')
            return text.strip()
            
        except Exception as e:
            logger.warning(f"Image OCR failed: {e}")
            return ""
    
    async def parse_pptx(self, file_path: Path) -> ParsedContent:
        """Parse PowerPoint PPTX files"""
        try:
            from pptx import Presentation
            
            logger.info(f"Parsing PPTX: {file_path}")
            prs = Presentation(str(file_path))
            
            text_parts = []
            images = []
            tables = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"# Slide {slide_num}\n")
                
                for shape in slide.shapes:
                    # Extract text
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
                    
                    # Extract tables
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        
                        if table_data:
                            tables.append({
                                "page": slide_num,
                                "table_index": len(tables),
                                "data": table_data,
                                "caption": f"Table on slide {slide_num}"
                            })
                            table_text = self._format_table_as_text(table_data)
                            text_parts.append(f"\n[Table]\n{table_text}\n")
                    
                    # Extract images
                    if shape.shape_type == 13:  # Picture
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            image_b64 = base64.b64encode(image_bytes).decode('utf-8')
                            ocr_text = ""
                            if self._should_ocr_embedded_image(image_bytes, ocr_count):
                                logger.info(
                                    f"[AdvancedParser] PPTX OCR image on slide {slide_num} ({len(image_bytes)//1024}KB)"
                                )
                                ocr_text = await self._ocr_image(image_bytes)
                                ocr_count += 1
                            
                            images.append({
                                "page": slide_num,
                                "image_index": len(images),
                                "format": image.ext,
                                "data": image_b64,
                                "caption": f"Image on slide {slide_num}",
                                "extracted_text": ocr_text
                            })
                            
                            if ocr_text:
                                text_parts.append(f"\n[Image Text: {ocr_text}]\n")
                        except Exception as e:
                            logger.warning(f"Error extracting image from slide {slide_num}: {e}")
                
                text_parts.append("\n")
            
            full_text = "\n".join(text_parts)
            
            metadata = {
                "total_slides": len(prs.slides),
                "core_properties": {
                    "title": prs.core_properties.title or "",
                    "author": prs.core_properties.author or "",
                    "subject": prs.core_properties.subject or "",
                }
            }
            if ocr_count > 0:
                logger.info(
                    f"[AdvancedParser] PPTX OCR'd {ocr_count} images "
                    f"(skipped {max(0, len(images) - ocr_count)} small/excess)"
                )
            
            logger.info(f"✅ Parsed PPTX: {len(full_text)} chars, {len(prs.slides)} slides")
            
            return ParsedContent(
                text=full_text,
                images=images,
                tables=tables,
                metadata=metadata,
                total_pages=len(prs.slides)
            )
            
        except ImportError:
            raise ValueError("python-pptx not installed. Cần rebuild Docker image.")
        except Exception as e:
            logger.error(f"PPTX parsing error: {e}")
            raise ValueError(f"Lỗi parse PPTX: {str(e)}")
    
    async def parse_xlsx(self, file_path: Path) -> ParsedContent:
        """Parse Excel XLSX files"""
        try:
            import openpyxl
            
            logger.info(f"Parsing XLSX: {file_path}")
            wb = openpyxl.load_workbook(str(file_path), data_only=True)
            
            text_parts = []
            tables = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"# Sheet: {sheet_name}\n")
                
                # Extract as table
                table_data = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_data = [str(cell) if cell is not None else "" for cell in row]
                        table_data.append(row_data)
                
                if table_data:
                    tables.append({
                        "page": 0,
                        "table_index": len(tables),
                        "data": table_data,
                        "caption": f"Sheet: {sheet_name}"
                    })
                    
                    table_text = self._format_table_as_text(table_data)
                    text_parts.append(f"{table_text}\n\n")
            
            full_text = "\n".join(text_parts)
            
            metadata = {
                "total_sheets": len(wb.sheetnames),
                "sheet_names": wb.sheetnames,
                "properties": {
                    "title": wb.properties.title or "",
                    "creator": wb.properties.creator or "",
                }
            }
            
            logger.info(f"✅ Parsed XLSX: {len(full_text)} chars, {len(wb.sheetnames)} sheets")
            
            return ParsedContent(
                text=full_text,
                images=[],
                tables=tables,
                metadata=metadata,
                total_pages=len(wb.sheetnames)
            )
            
        except ImportError:
            raise ValueError("openpyxl not installed. Cần rebuild Docker image.")
        except Exception as e:
            logger.error(f"XLSX parsing error: {e}")
            raise ValueError(f"Lỗi parse XLSX: {str(e)}")
    
    async def parse_markdown(self, file_path: Path) -> ParsedContent:
        """Parse Markdown files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            logger.info(f"✅ Parsed Markdown: {len(content)} chars")
            
            return ParsedContent(
                text=content,
                images=[],
                tables=[],
                metadata={"format": "markdown"},
                total_pages=1
            )
        except Exception as e:
            logger.error(f"Markdown parsing error: {e}")
            raise ValueError(f"Lỗi parse Markdown: {str(e)}")
    
    async def parse_html(self, file_path: Path) -> ParsedContent:
        """Parse HTML files"""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Extract text
            text = soup.get_text(separator='\n', strip=True)
            
            # Extract tables
            tables = []
            for table in soup.find_all('table'):
                table_data = []
                for row in table.find_all('tr'):
                    row_data = [cell.get_text(strip=True) for cell in row.find_all(['td', 'th'])]
                    if row_data:
                        table_data.append(row_data)
                
                if table_data:
                    tables.append({
                        "page": 0,
                        "table_index": len(tables),
                        "data": table_data,
                        "caption": f"HTML Table {len(tables) + 1}"
                    })
            
            logger.info(f"✅ Parsed HTML: {len(text)} chars, {len(tables)} tables")
            
            return ParsedContent(
                text=text,
                images=[],
                tables=tables,
                metadata={"format": "html", "title": soup.title.string if soup.title else ""},
                total_pages=1
            )
            
        except ImportError:
            raise ValueError("beautifulsoup4 not installed. Cần rebuild Docker image.")
        except Exception as e:
            logger.error(f"HTML parsing error: {e}")
            raise ValueError(f"Lỗi parse HTML: {str(e)}")
    
    async def parse_rtf(self, file_path: Path) -> ParsedContent:
        """Parse RTF files"""
        try:
            from striprtf.striprtf import rtf_to_text
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            
            text = rtf_to_text(rtf_content)
            
            logger.info(f"✅ Parsed RTF: {len(text)} chars")
            
            return ParsedContent(
                text=text,
                images=[],
                tables=[],
                metadata={"format": "rtf"},
                total_pages=1
            )
            
        except ImportError:
            raise ValueError("striprtf not installed. Cần rebuild Docker image.")
        except Exception as e:
            logger.error(f"RTF parsing error: {e}")
            raise ValueError(f"Lỗi parse RTF: {str(e)}")
    
    async def parse_odt(self, file_path: Path) -> ParsedContent:
        """Parse OpenDocument Text files"""
        try:
            from odf import text, teletype
            from odf.opendocument import load
            
            doc = load(str(file_path))
            text_parts = []
            
            for paragraph in doc.getElementsByType(text.P):
                text_parts.append(teletype.extractText(paragraph))
            
            full_text = "\n".join(text_parts)
            
            logger.info(f"✅ Parsed ODT: {len(full_text)} chars")
            
            return ParsedContent(
                text=full_text,
                images=[],
                tables=[],
                metadata={"format": "odt"},
                total_pages=1
            )
            
        except ImportError:
            raise ValueError("odfpy not installed. Cần rebuild Docker image.")
        except Exception as e:
            logger.error(f"ODT parsing error: {e}")
            raise ValueError(f"Lỗi parse ODT: {str(e)}")
    
    async def parse_csv(self, file_path: Path) -> ParsedContent:
        """Parse CSV files"""
        try:
            import csv
            
            with open(file_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                table_data = list(reader)
            
            # Format as text
            text_parts = []
            for row in table_data:
                text_parts.append(" | ".join(row))
            
            full_text = "\n".join(text_parts)
            
            tables = [{
                "page": 0,
                "table_index": 0,
                "data": table_data,
                "caption": "CSV Data"
            }]
            
            logger.info(f"✅ Parsed CSV: {len(table_data)} rows")
            
            return ParsedContent(
                text=full_text,
                images=[],
                tables=tables,
                metadata={"format": "csv", "rows": len(table_data)},
                total_pages=1
            )
            
        except Exception as e:
            logger.error(f"CSV parsing error: {e}")
            raise ValueError(f"Lỗi parse CSV: {str(e)}")
    
    async def parse_json(self, file_path: Path) -> ParsedContent:
        """Parse JSON files"""
        try:
            import json
            
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Convert to readable text
            text = json.dumps(data, indent=2, ensure_ascii=False)
            
            logger.info(f"✅ Parsed JSON: {len(text)} chars")
            
            return ParsedContent(
                text=text,
                images=[],
                tables=[],
                metadata={"format": "json"},
                total_pages=1
            )
            
        except Exception as e:
            logger.error(f"JSON parsing error: {e}")
            raise ValueError(f"Lỗi parse JSON: {str(e)}")
    
    async def parse_image(self, file_path: Path) -> ParsedContent:
        """Parse image files (pure OCR)"""
        try:
            with open(file_path, 'rb') as f:
                image_bytes = f.read()
            
            # OCR the image
            ocr_text = await self._ocr_image(image_bytes)
            
            image_b64 = base64.b64encode(image_bytes).decode('utf-8')
            
            images = [{
                "page": 1,
                "image_index": 0,
                "format": file_path.suffix[1:],
                "data": image_b64,
                "caption": file_path.name,
                "extracted_text": ocr_text
            }]
            
            logger.info(f"✅ Parsed Image: {len(ocr_text)} chars OCR text")
            
            return ParsedContent(
                text=ocr_text or "[Image with no extractable text]",
                images=images,
                tables=[],
                metadata={"format": "image", "filename": file_path.name},
                total_pages=1
            )
            
        except Exception as e:
            logger.error(f"Image parsing error: {e}")
            raise ValueError(f"Lỗi parse Image: {str(e)}")
    
    def _format_table_as_text(self, table_data: List[List[str]]) -> str:
        """Format table as markdown-style text"""
        if not table_data:
            return ""
        
        lines = []
        for row in table_data:
            lines.append(" | ".join(str(cell) for cell in row))
        
        return "\n".join(lines)


# Singleton
_parser_instance = None


def get_advanced_parser() -> AdvancedDocumentParser:
    global _parser_instance
    if _parser_instance is None:
        _parser_instance = AdvancedDocumentParser()
    return _parser_instance
